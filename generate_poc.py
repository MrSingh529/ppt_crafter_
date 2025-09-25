import openpyxl
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from copy import deepcopy
import math
from pptx.oxml.ns import qn
import requests, re, json, datetime
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Emu
import google.generativeai as genai

# setup Gemini
genai.configure(api_key="AIzaSyD4LOjhJtiC2mdRZJYGSa5iavCBRTZfTKU")

# --------------- helpers ---------------

def read_summary_keys(excel_path, sheet_name="Summary"):
    wb = openpyxl.load_workbook(excel_path, data_only=True)
    ws = wb[sheet_name]
    kv = {}
    for row in ws.iter_rows(min_row=2, max_row=ws.max_row,
                            min_col=1, max_col=2, values_only=True):
        key, val = row
        if key is None:
            continue
        kv[str(key).strip()] = "" if val is None else str(val)
    return kv

def find_header_row(ws, search_year=2024, percent=False):
    """Find the row index that contains the given year.
       If percent=True, only match headers with '(%)' in them."""
    for i, row in enumerate(ws.iter_rows(min_row=1, max_row=40, values_only=True), start=1):
        values = [str(c) if c else "" for c in row]
        if any(str(search_year) in v for v in values):
            if not percent or any("(%)" in v for v in values):
                return i, values
    raise ValueError(f"Header row with {search_year}{' (%)' if percent else ''} not found in {ws.title}")

def _is_paragraph_bulleted(paragraph):
    pPr = getattr(paragraph._p, "pPr", None)
    if pPr is None: 
        return False
    if pPr.find(qn("a:buNone")) is not None:
        return False
    return any(pPr.find(qn(tag)) is not None for tag in ("a:buChar","a:buAutoNum","a:buBlip"))

def _pick_ref_paragraph(text_frame, target_para):
    # Prefer an existing bulleted paragraph in the same text frame
    if _is_paragraph_bulleted(target_para):
        return target_para
    for para in text_frame.paragraphs:
        if para is not target_para and _is_paragraph_bulleted(para):
            return para
    return target_para

def _clone_pPr(dst_para, src_pPr):
    # Replace dst_para's pPr with a clone of src_pPr using lxml API (no _element)
    _ = dst_para._p.get_or_add_pPr()
    try:
        if getattr(dst_para._p, "pPr", None) is not None:
            dst_para._p.remove(dst_para._p.pPr)
    except Exception:
        pass
    dst_para._p.insert(0, deepcopy(src_pPr))

def find_year_col(header, year):
    """Find column index where header contains the given year."""
    year_str = str(year)
    for idx, h in enumerate(header):
        if h and year_str in str(h):
            return idx
    raise ValueError(f"Year {year} not found in headers: {header}")

def generate_overview_ai_content(excel_path, existing_kv=None, use_ai=True):
    """
    Generate detailed overview content using AI based on Excel data and market information
    """
    if not use_ai:
        return ""
    
    # Use provided kv data if available, otherwise extract fresh
    if existing_kv:
        kv = existing_kv
    else:
        kv = read_summary_keys(excel_path, "Summary")
        # Extract basic dynamic data for AI prompt
        wb = openpyxl.load_workbook(excel_path, data_only=True)
        
        title = wb["Summary"]["B2"].value if "Summary" in wb.sheetnames else ""
        if title:
            parts = title.split()
            country = parts[0] if len(parts) > 0 else ""
            product = parts[1] if len(parts) > 1 else ""
        else:
            country = ""
            product = ""
        
        units = wb["Summary"]["B4"].value if "Summary" in wb.sheetnames else ""
        unit = "Thousand Tons" if units and "Thousand Tons" in units else units or ""
        
        kv.update({
            "Title": title,
            "Country": country,
            "Product": product,
            "Unit": unit,
        })
    
    # Extract current market size and key data points
    wb = openpyxl.load_workbook(excel_path, data_only=True)
    
    # Get basic market data for context
    market_size = kv.get('Sales_Volume_Latest', '')
    latest_year = kv.get('Latest_Year', '2024')
    
    try:
        prompt = f"""
        Write a detailed and exhaustive overview for the {kv.get('Title', '')} market. Use paragraph form and write exactly 350 words.
        
        Market Context:
        - Market: {kv.get('Product', '')} in {kv.get('Country', '')}
        - Current Market Size ({latest_year}): {market_size} {kv.get('Unit', '')}
        - Historical Growth: {kv.get('CAGR_2019_2024', '')} CAGR
        - Forecast Growth: {kv.get('CAGR_2025_2033', '')} CAGR
        
        Requirements:
        1. Write in paragraph form (no bullet points or lists)
        2. Exactly 350 words
        3. Cover the chemical compound's basic properties and characteristics
        4. Include historical context and evolution of the market
        5. Discuss key applications and industrial uses
        6. Mention production processes and technological aspects
        7. Address market advantages and benefits
        8. Include regulatory and environmental considerations
        9. Discuss market drivers and growth factors
        10. Provide comprehensive technical and commercial overview
        
        Write a comprehensive, technical, and market-focused overview that would be suitable for an industry report introduction section.
        """
        
        model = genai.GenerativeModel("gemini-1.5-flash")
        response = model.generate_content(prompt)
        
        content = getattr(response, "text", "").strip()
        
        # Clean up any unwanted formatting
        content = re.sub(r'\*\*([^*]+)\*\*', r'\1', content)  # Remove bold markdown
        content = re.sub(r'\*([^*]+)\*', r'\1', content)      # Remove italic markdown
        content = re.sub(r'#+\s*', '', content)              # Remove headers
        content = re.sub(r'^\s*[-•]\s*', '', content, flags=re.MULTILINE)  # Remove bullet points
        
        return content
        
    except Exception as e:
        print(f"AI overview content generation failed: {e}")
        return ""

def generate_market_overview_content(excel_path, existing_kv=None, use_ai=True):
    """
    Generate detailed market overview content using AI based on Excel data
    Takes existing_kv to avoid circular dependency
    """
    if not use_ai:
        return ""
    
    # Use provided kv data if available, otherwise extract fresh (but don't include Market_Overview_Content)
    if existing_kv:
        kv = existing_kv
    else:
        kv = read_summary_keys(excel_path, "Summary")
        # Extract basic dynamic data WITHOUT calling the full extract_dynamic_placeholders
        wb = openpyxl.load_workbook(excel_path, data_only=True)
        
        # Extract only essential data needed for AI prompt
        title = wb["Summary"]["B2"].value if "Summary" in wb.sheetnames else ""
        if title:
            parts = title.split()
            country = parts[0] if len(parts) > 0 else ""
            product = parts[1] if len(parts) > 1 else ""
        else:
            country = ""
            product = ""
        
        units = wb["Summary"]["B4"].value if "Summary" in wb.sheetnames else ""
        unit = "Thousand Tons" if units and "Thousand Tons" in units else units or ""
        
        # Sales Forecast data
        if "Sales_Forecast" in wb.sheetnames:
            ws = wb["Sales_Forecast"]
            years = [r[0] for r in ws.iter_rows(min_row=2, values_only=True) if r[0]]
            volumes = {row[0]: row[1] for row in ws.iter_rows(min_row=2, values_only=True)}
            latest_year = max(y for y in years if y <= 2024) if years else 2024
            
            kv.update({
                "Title": title,
                "Country": country,
                "Product": product,
                "Unit": unit,
                "Latest_Year": str(latest_year),
                "Sales_Volume_Latest": f"{volumes.get(latest_year, 0):,.0f}",
                "Sales_Volume_2033": f"{volumes.get(2033, 0):,.0f}",
            })
            
            # Get CAGR values
            def fmt_pct(val):
                try:
                    return f"{float(val)*100:.1f}%" if abs(float(val)) < 1 else f"{float(val):.1f}%"
                except:
                    return str(val) if val else ""
            
            kv["CAGR_2019_2024"] = fmt_pct(ws["C7"].value) if ws["C7"].value else ""
            kv["CAGR_2025_2033"] = fmt_pct(ws["D16"].value) if ws["D16"].value else ""
    
    # Extract segmentation data for AI prompt
    wb = openpyxl.load_workbook(excel_path, data_only=True)
    
    # Get top segments from each sheet
    type_data = get_sheet_percentage_data("By_Type", wb)
    app_data = get_sheet_percentage_data("By_Application", wb)
    enduser_data = get_sheet_percentage_data("By_EndUser", wb)
    region_data = get_sheet_percentage_data("By_Region", wb)
    
    try:
        prompt = f"""
        Write a detailed and exhaustive market overview for the {kv.get('Title', '')} market in exactly 230 words. Use paragraph form.
        
        Key Market Data:
        - Market: {kv.get('Product', '')} in {kv.get('Country', '')}
        - Current Volume ({kv.get('Latest_Year', '')}): {kv.get('Sales_Volume_Latest', '')} {kv.get('Unit', '')}
        - Projected Volume (2033): {kv.get('Sales_Volume_2033', '')} {kv.get('Unit', '')}  
        - Historical CAGR (2019-2024): {kv.get('CAGR_2019_2024', '')}
        - Forecast CAGR (2025-2033): {kv.get('CAGR_2025_2033', '')}
        
        Market Segmentation:
        - Top Product Types: {', '.join([f"{name} ({share}%)" for name, share in type_data[:3]])}
        - Key Applications: {', '.join([f"{name} ({share}%)" for name, share in app_data[:3]])}
        - Major End Users: {', '.join([f"{name} ({share}%)" for name, share in enduser_data[:3]])}
        - Regional Distribution: {', '.join([f"{name} ({share}%)" for name, share in region_data[:3]])}
        
        Requirements:
        1. Write in paragraph form (no bullet points)
        2. Exactly 500 words
        3. Include market drivers, challenges, opportunities
        4. Mention key segments and their significance  
        5. Discuss technological trends and innovations
        6. Cover regulatory environment and policy impacts
        7. Address competitive landscape
        8. Include future outlook and growth prospects
        
        Focus on providing comprehensive market intelligence that would be valuable for business decision-making.
        """
        
        model = genai.GenerativeModel("gemini-1.5-flash")
        response = model.generate_content(prompt)
        
        content = getattr(response, "text", "").strip()
        
        # Clean up any unwanted formatting
        content = re.sub(r'\*\*([^*]+)\*\*', r'\1', content)  # Remove bold markdown
        content = re.sub(r'\*([^*]+)\*', r'\1', content)      # Remove italic markdown
        content = re.sub(r'#+\s*', '', content)              # Remove headers
        
        return content
        
    except Exception as e:
        print(f"AI content generation failed: {e}")
        return ""

def get_sheet_percentage_data(sheet_name, workbook):
    """Extract percentage data from segmentation sheets - simplified version"""
    if sheet_name not in workbook.sheetnames:
        return []
    
    ws = workbook[sheet_name]
    data = []
    
    # Look for percentage section
    for row_idx in range(1, 20):
        row = list(ws.iter_rows(min_row=row_idx, max_row=row_idx, values_only=True))[0]
        row_text = " ".join(str(cell) if cell else "" for cell in row)
        
        if "(%" in row_text and "2024" in row_text:
            # Found percentage section, extract data
            for data_row_idx in range(row_idx + 1, min(row_idx + 10, ws.max_row + 1)):
                data_row = list(ws.iter_rows(min_row=data_row_idx, max_row=data_row_idx, values_only=True))[0]
                
                if not data_row[0] or "Total" in str(data_row[0]):
                    break
                
                item_name = str(data_row[0]).strip()
                try:
                    # Find the percentage column (usually column 1 or 2)
                    value = None
                    for col_idx in range(1, min(3, len(data_row))):
                        if data_row[col_idx] is not None:
                            value = float(data_row[col_idx])
                            if 0 < value < 1:
                                value = round(value * 100, 1)
                            elif value:
                                value = round(value, 1)
                            break
                    
                    if value is not None:
                        data.append((item_name, value))
                except (ValueError, TypeError, IndexError):
                    continue
            break
    
    return sorted(data, key=lambda x: x[1] if x[1] is not None else 0, reverse=True)

def extract_dynamic_placeholders(excel_path, include_market_overview=True, include_overview_content=True):
    """Extract dynamic placeholders from Sales_Forecast + segmentation sheets."""
    wb = openpyxl.load_workbook(excel_path, data_only=True)
    kv = {}

    # --- Title ---
    title = wb["Summary"]["B2"].value
    kv["Title"] = title
    if title:
        parts = title.split()
        kv["Country"] = parts[0] if len(parts) > 0 else ""
        kv["Product"] = parts[1] if len(parts) > 1 else ""

    # --- Units ---
    units = wb["Summary"]["B4"].value if "Summary" in wb.sheetnames else ""
    kv["Unit"] = "Thousand Tons" if units and "Thousand Tons" in units else units

    # --- Sales Forecast ---
    ws = wb["Sales_Forecast"]
    years = [r[0] for r in ws.iter_rows(min_row=2, values_only=True) if r[0]]
    volumes = {int(row[0]): float(row[1]) if row[1] is not None else 0.0 for row in ws.iter_rows(min_row=2, values_only=True) if row[0] is not None}
    latest_year = max(y for y in years if y <= 2024)
    kv["Latest_Year"] = str(latest_year)
    kv["Sales_Volume_Latest"] = f"{volumes.get(int(latest_year), 0):,.0f}"
    kv["Sales_Volume_2033"] = f"{volumes.get(2033, 0):,.0f}"
    kv["Historical_Start_Year"] = "2019"
    kv["Historical_End_Year"] = str(latest_year)
    kv["Forecast_Start_Year"] = "2025"
    kv["Forecast_End_Year"] = "2033"

    # Format CAGRs nicely
    def fmt_pct(val):
        try:
            return f"{float(val)*100:.1f}%" if abs(float(val)) < 1 else f"{float(val):.1f}%"
        except:
            return str(val) if val else ""

    kv["CAGR_2019_2024"] = fmt_pct(ws["C7"].value)
    kv["CAGR_2025_2033"] = fmt_pct(ws["D16"].value)

    try:
        cagr_val = float(str(ws["C7"].value))
        kv["Trend_Phrase"] = "growing" if cagr_val > 0 else "declining" if cagr_val < 0 else "remaining stable"
    except:
        kv["Trend_Phrase"] = ""

    # --- Updated processor for volume data only ---
    def process_sheet_volume_data(sheet_name, top_n=None):
        """Process sheet to get VOLUME data (not percentage data)"""
        ws = wb[sheet_name]
        
        # Use the updated function to get volume data for latest year
        volume_data = get_sheet_data_for_year(excel_path, sheet_name, latest_year)
        
        # Convert to list of tuples and sort
        data = [(name, value) for name, value in volume_data.items()]
        data.sort(key=lambda x: x[1] if x[1] is not None else 0, reverse=True)
        
        return data if not top_n else data[:top_n], data

    # --- Updated processor for percentage data ---
    def process_sheet_percentage_data(sheet_name, top_n=None):
        """Process sheet to get PERCENTAGE data specifically"""
        ws = wb[sheet_name]
        
        # Look for percentage section
        header_row_idx = None
        year_col_idx = None
        
        for row_idx in range(1, 20):
            row = list(ws.iter_rows(min_row=row_idx, max_row=row_idx, values_only=True))[0]
            row_text = " ".join(str(cell) if cell else "" for cell in row)
            
            # Look for percentage section AND the year
            if "(%" in row_text and str(latest_year) in row_text:
                for col_idx, cell in enumerate(row):
                    if cell and str(latest_year) in str(cell) and "(%" in str(cell):
                        header_row_idx = row_idx
                        year_col_idx = col_idx
                        break
                break
        
        if header_row_idx is None or year_col_idx is None:
            return [], []
        
        # Extract percentage data
        data = []
        for row_idx in range(header_row_idx + 1, min(header_row_idx + 10, ws.max_row + 1)):
            row = list(ws.iter_rows(min_row=row_idx, max_row=row_idx, values_only=True))[0]
            
            if not row[0] or "Total" in str(row[0]):
                break
                
            item_name = str(row[0]).strip()
            try:
                value = float(row[year_col_idx]) if row[year_col_idx] is not None else 0
                # Convert to percentage if it's a decimal
                if 0 < value < 1:
                    value = round(value * 100, 1)
                elif value:
                    value = round(value, 1)
                data.append((item_name, value))
            except (ValueError, TypeError, IndexError):
                continue
        
        data.sort(key=lambda x: x[1] if x[1] is not None else 0, reverse=True)
        return data if not top_n else data[:top_n], data

    # By_Type (top 3 + aliases) - USE PERCENTAGE DATA
    type_top, _ = process_sheet_percentage_data("By_Type", top_n=None)
    for i, (name, val) in enumerate(type_top, start=1):
        kv[f"Top_Type_{i}"] = name
        kv[f"Top_Type_{i}_Share"] = f"{val:.1f}" if val is not None else ""
    if len(type_top) >= 1:
        kv["Top_Type"] = type_top[0][0]
        kv["Top_Type_Share"] = f"{type_top[0][1]:.1f}"
    if len(type_top) >= 2:
        kv["Second_Type"] = type_top[1][0]
        kv["Second_Type_Share"] = f"{type_top[1][1]:.1f}"
    if len(type_top) >= 3:
        kv["Third_Type"] = type_top[2][0]
        kv["Third_Type_Share"] = f"{type_top[2][1]:.1f}"

    # By_Application (top 5 + aliases + others) - USE PERCENTAGE DATA
    app_top, app_all = process_sheet_percentage_data("By_Application")
    for i, (name, val) in enumerate(app_top[:5], start=1):
        kv[f"Top_Application_{i}"] = name
        kv[f"Top_Application_{i}_Share"] = f"{val:.1f}" if val is not None else ""
    aliases = ["Top", "Second", "Third", "Fourth", "Fifth"]
    for i, alias in enumerate(aliases, start=1):
        if len(app_top) >= i:
            kv[f"{alias}_Application"] = app_top[i-1][0]
            kv[f"{alias}_Application_Share"] = f"{app_top[i-1][1]:.1f}"
    if len(app_all) > 5:
        kv["Other_Application_Share"] = f"{sum(v for _, v in app_all[5:] if v):.1f}"

    # By_EndUser (top 5 + aliases + others) - USE PERCENTAGE DATA
    eu_top, eu_all = process_sheet_percentage_data("By_EndUser")
    for i, (name, val) in enumerate(eu_top[:5], start=1):
        kv[f"Top_EndUser_{i}"] = name
        kv[f"Top_EndUser_{i}_Share"] = f"{val:.1f}" if val is not None else ""
    aliases = ["Top", "Second", "Third", "Fourth", "Fifth"]
    for i, alias in enumerate(aliases, start=1):
        if len(eu_top) >= i:
            kv[f"{alias}_EndUser"] = eu_top[i-1][0]
            kv[f"{alias}_EndUser_Share"] = f"{eu_top[i-1][1]:.1f}"
    if len(eu_all) > 5:
        kv["Other_EndUser_Share"] = f"{sum(v for _, v in eu_all[5:] if v):.1f}"

    # By_Region (all dynamically, from % block) - USE PERCENTAGE DATA
    reg_all, _ = process_sheet_percentage_data("By_Region")
    for i, (name, val) in enumerate(reg_all, start=1):
        kv[f"Top_Region_{i}"] = name
        kv[f"Top_Region_{i}_Share"] = f"{val:.1f}" if val is not None else ""

    # --- Narrative lines ---
    kv["Market_Intro_Line"] = (
        f"The {kv['Product'].lower()} market in {kv['Country']} reached a volume of "
        f"{kv['Sales_Volume_Latest']} {kv['Unit']} in {kv['Latest_Year']}, "
        f"{kv['Trend_Phrase']} at a CAGR of {kv['CAGR_2019_2024']} during 2019–2024."
    )

    kv["Market_Outlook_Line"] = (
        f"Overall, the {kv['Product'].lower()} market in {kv['Country']} is expected to grow "
        f"at a CAGR of {kv['CAGR_2025_2033']} during 2025–2033, reaching sales worth "
        f"{kv['Sales_Volume_2033']} {kv['Unit']} by 2033."
    )

    if include_market_overview:
        # Generate market overview content using the data we just extracted
        kv["Market_Overview_Content"] = generate_market_overview_content(excel_path, existing_kv=kv, use_ai=True)
    
    # NEW: Generate overview AI content
    if include_overview_content:
        kv["Overview_AI_Content"] = generate_overview_ai_content(excel_path, existing_kv=kv, use_ai=True)
    
    return kv, volumes

def build_report_subtitle(excel_path):
    wb = openpyxl.load_workbook(excel_path, data_only=True)

    def collect_items(sheet_name):
        ws = wb[sheet_name]
        items = []
        for row in ws.iter_rows(min_row=2, values_only=True):
            val = row[0]
            if not val:
                continue
            val_str = str(val).strip()
            low = val_str.lower()
            # filter out headers / totals
            if (
                low.startswith("type") or low.startswith("source") or low.startswith("end user")
                or low.startswith("region") or low.startswith("total")
                or "market breakup" in low
            ):
                continue
            items.append(val_str)
        return items

    type_items = collect_items("By_Type")
    app_items = collect_items("By_Application")
    eu_items = collect_items("By_EndUser")

    # reorder Application items to match example (move "Others" to last)
    if "Others" in app_items:
        app_items = [i for i in app_items if i != "Others"] + ["Others"]

    # reorder EndUser items similarly
    if "Others" in eu_items:
        eu_items = [i for i in eu_items if i != "Others"] + ["Others"]

    subtitle = (
        f"Report by Physical Form ({', '.join(type_items)}), "
        f"Application ({', '.join(app_items)}), "
        f"End Use Industry ({', '.join(eu_items)}), "
        f"and Region 2025–2033"
    )
    return subtitle

def build_list_from_sheet(excel_path, sheet_name, ignore_headers=True):
    wb = openpyxl.load_workbook(excel_path, data_only=True)
    ws = wb[sheet_name]
    items = []
    seen = set()
    for row in ws.iter_rows(min_row=2, max_row=ws.max_row, min_col=1, max_col=1, values_only=True):
        val = row[0]
        if not val:
            continue
        val = str(val).strip()
        low = val.lower()
        if ignore_headers and (
            low.startswith("type") or low.startswith("source") or low.startswith("end user")
            or low.startswith("region") or low.startswith("total") or low.startswith("market breakup")
        ):
            continue
        if val in seen:
            continue
        seen.add(val)
        items.append(val)
    return items

# NEW FUNCTION: Create inline text versions of lists
def create_inline_placeholders(excel_path):
    """Create inline (comma-separated) versions of list placeholders."""
    wb = openpyxl.load_workbook(excel_path, data_only=True)
    inline_kv = {}
    
    # Define the sheets we want to create inline versions for
    sheet_mappings = {
        "By_Type": "By_Type_Inline",
        "By_Application": "By_Application_Inline", 
        "By_EndUser": "By_EndUser_Inline",
        "By_Region": "By_Region_Inline"
    }
    
    for sheet_name, inline_key in sheet_mappings.items():
        if sheet_name in wb.sheetnames:
            items = build_list_from_sheet(excel_path, sheet_name)
            # Create comma-separated inline text
            inline_text = ", ".join(items)
            inline_kv[inline_key] = inline_text
        else:
            inline_kv[inline_key] = ""
    
    return inline_kv

def build_toc_from_sheet(excel_path, sheet_name="Table_Contents"):
    """Return list of (text, level) from Table_Contents sheet."""
    wb = openpyxl.load_workbook(excel_path, data_only=True)
    ws = wb[sheet_name]
    toc_items = []
    for row in ws.iter_rows(min_row=1, max_row=ws.max_row, min_col=1, max_col=1, values_only=True):
        val = row[0]
        if not val:
            continue
        text = str(val).strip()
        # Level = count of dots in prefix numbering
        parts = text.split(" ", 1)
        if len(parts) > 1 and parts[0][0].isdigit():
            numbering = parts[0]
            level = numbering.count(".")  # 0=digit only, 1=one dot, etc.
        else:
            level = 0
        toc_items.append((text, level))
    return toc_items

def safe_copy_font(src_font, dst_font):
    """Improved font copying that better preserves all attributes"""
    try:
        # Copy name
        if hasattr(src_font, 'name') and src_font.name:
            dst_font.name = src_font.name
        
        # Copy size - be more careful about None values
        if hasattr(src_font, 'size') and src_font.size is not None:
            dst_font.size = src_font.size
        
        # Copy bold
        if hasattr(src_font, 'bold') and src_font.bold is not None:
            dst_font.bold = src_font.bold
        
        # Copy italic  
        if hasattr(src_font, 'italic') and src_font.italic is not None:
            dst_font.italic = src_font.italic
        
        # Copy underline
        if hasattr(src_font, 'underline') and src_font.underline is not None:
            dst_font.underline = src_font.underline
        
        # Copy color more carefully
        if hasattr(src_font, 'color') and src_font.color:
            try:
                if hasattr(src_font.color, 'rgb') and src_font.color.rgb:
                    dst_font.color.rgb = src_font.color.rgb
                elif hasattr(src_font.color, 'theme_color') and src_font.color.theme_color is not None:
                    # Only set valid theme colors
                    from pptx.enum.dml import MSO_THEME_COLOR_INDEX
                    if src_font.color.theme_color != MSO_THEME_COLOR_INDEX.NOT_THEME_COLOR:
                        dst_font.color.theme_color = src_font.color.theme_color
            except Exception as e:
                print(f"Warning: Could not copy color: {e}")
                
    except Exception as e:
        print(f"Warning: Error copying font properties: {e}")

def get_run_formatting(run):
    """Extract comprehensive formatting from a run"""
    formatting = {}
    
    if not run or not hasattr(run, 'font'):
        return formatting
        
    font = run.font
    
    # Extract all font properties safely
    try:
        formatting['font_name'] = getattr(font, 'name', None)
        formatting['font_size'] = getattr(font, 'size', None)
        formatting['font_bold'] = getattr(font, 'bold', None)
        formatting['font_italic'] = getattr(font, 'italic', None)
        formatting['font_underline'] = getattr(font, 'underline', None)
        
        # Handle color
        if hasattr(font, 'color') and font.color:
            if hasattr(font.color, 'rgb') and font.color.rgb:
                formatting['font_color_rgb'] = font.color.rgb
            elif hasattr(font.color, 'theme_color') and font.color.theme_color is not None:
                from pptx.enum.dml import MSO_THEME_COLOR_INDEX
                if font.color.theme_color != MSO_THEME_COLOR_INDEX.NOT_THEME_COLOR:
                    formatting['font_color_theme'] = font.color.theme_color
    except Exception as e:
        print(f"Warning: Error extracting run formatting: {e}")
    
    return formatting

def apply_run_formatting(run, formatting):
    """Apply formatting to a run"""
    if not formatting or not run or not hasattr(run, 'font'):
        return
    
    font = run.font
    
    try:
        # Apply font properties only if they exist in formatting
        if formatting.get('font_name'):
            font.name = formatting['font_name']
        if formatting.get('font_size') is not None:
            font.size = formatting['font_size']
        if formatting.get('font_bold') is not None:
            font.bold = formatting['font_bold']
        if formatting.get('font_italic') is not None:
            font.italic = formatting['font_italic']
        if formatting.get('font_underline') is not None:
            font.underline = formatting['font_underline']
        
        # Apply color - RGB takes priority
        if formatting.get('font_color_rgb'):
            font.color.rgb = formatting['font_color_rgb']
        elif formatting.get('font_color_theme') is not None:
            font.color.theme_color = formatting['font_color_theme']
    except Exception as e:
        print(f"Warning: Error applying run formatting: {e}")

def copy_shape_style(src_shape, dst_shape):
    """Deep copy style: fill, border, margins, alignment, font, paragraph spacing"""
    try:
        if src_shape.fill and src_shape.fill.type is not None:
            dst_shape.fill.solid()
            dst_shape.fill.fore_color.rgb = src_shape.fill.fore_color.rgb
    except:
        pass

    try:
        if src_shape.line and src_shape.line.color and src_shape.line.color.rgb:
            dst_shape.line.color.rgb = src_shape.line.color.rgb
    except:
        pass

    if src_shape.has_text_frame and dst_shape.has_text_frame:
        src_tf, dst_tf = src_shape.text_frame, dst_shape.text_frame

        # Copy margins + wrap
        dst_tf.margin_left = src_tf.margin_left
        dst_tf.margin_right = src_tf.margin_right
        dst_tf.margin_top = src_tf.margin_top
        dst_tf.margin_bottom = src_tf.margin_bottom
        dst_tf.word_wrap = src_tf.word_wrap

        # Copy paragraphs alignment + spacing
        if src_tf.paragraphs:
            for i, sp in enumerate(src_tf.paragraphs):
                if i >= len(dst_tf.paragraphs):
                    dst_tf.add_paragraph()
                dp = dst_tf.paragraphs[i]
                dp.alignment = sp.alignment
                dp.level = sp.level
                dp.space_after = sp.space_after
                dp.space_before = sp.space_before
                dp.line_spacing = sp.line_spacing

                # Copy font of first run
                if sp.runs:
                    src_font = sp.runs[0].font
                    if dp.runs:
                        safe_copy_font(src_font, dp.runs[0].font)

def add_extra_boxes(slide, template_shape, items, start_index=0, vertical_spacing=Emu(500000)):
    """
    Dynamically add extra shapes (same formatting) if items exceed placeholders.
    vertical_spacing default = ~0.5 inch
    """
    base_top = template_shape.top
    for i, item in enumerate(items[start_index:], start=0):
        new_el = deepcopy(template_shape.element)
        slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
        new_shape = slide.shapes[-1]

        # Position relative to template
        new_shape.left = template_shape.left
        new_shape.top = base_top + vertical_spacing * (i + 1)
        new_shape.width = template_shape.width
        new_shape.height = template_shape.height

        set_text_with_placeholder_format(new_shape, str(item))

# ----------------- PPT Modifiers -----------------

def duplicate_slide(prs, slide):
    """
    Duplicate a slide with layout & placeholders intact.
    Keeps placeholders so we can replace them later.
    """
    slide_layout = slide.slide_layout
    new_slide = prs.slides.add_slide(slide_layout)

    for shape in slide.shapes:
        new_el = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    return new_slide

def chunk_toc_items(toc_items, items_per_column=45):
    """
    Split TOC items into chunks for multiple slides.
    Each slide has two columns â†’ about 2 * items_per_column items max.
    """
    items_per_slide = items_per_column * 2
    for i in range(0, len(toc_items), items_per_slide):
        yield toc_items[i:i + items_per_slide]

def insert_toc_into_textframe(text_frame, toc_items, template_para=None):
    """Insert TOC into a given textframe with indent by level, no overflow handling here."""
    tmpl_para = template_para if template_para else (text_frame.paragraphs[0] if text_frame.paragraphs else None)
    tmpl_run = tmpl_para.runs[0] if tmpl_para and tmpl_para.runs else None
    tmpl_align = getattr(tmpl_para, "alignment", None)

    text_frame.text = ""
    for i, (text, level) in enumerate(toc_items):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = text
        p.level = level
        if tmpl_align is not None:
            p.alignment = tmpl_align
        if tmpl_run and p.runs:
            safe_copy_font(tmpl_run.font, p.runs[0].font)
        if level == 0 and p.runs:
            p.runs[0].font.bold = True

def estimate_items_per_column(shape, font_size_pt=14, line_spacing=1.2):
    """
    Estimate how many items fit in this shape's text_frame height.
    font_size_pt: average font size (pt)
    line_spacing: spacing multiplier
    """
    # shape.height is in EMUs, convert to inches
    frame_height_inch = shape.height / 914400.0  # 914400 EMUs = 1 inch
    line_height_inch = (font_size_pt / 72.0) * line_spacing  # 72pt = 1 inch
    max_items = int(frame_height_inch / line_height_inch)
    return max_items if max_items > 0 else 25  # fallback

def replace_toc_in_slide(slide, toc_items, items_per_column=30):
    """Replace TOC placeholders (left/right) in one slide and return leftovers."""
    left_tf, right_tf, left_para, right_para = None, None, None, None

    # find placeholders
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        tf = shape.text_frame
        for para in tf.paragraphs:
            txt = ''.join(r.text for r in para.runs)
            if "{{Table_Contents_Left}}" in txt:
                left_tf, left_para = tf, para
            elif "{{Table_Contents_Right}}" in txt:
                right_tf, right_para = tf, para

    if not left_tf or not right_tf:
        return toc_items  # nothing replaced, return everything

    # fixed split per column
    left = toc_items[:items_per_column]
    right = toc_items[items_per_column:items_per_column*2]
    leftovers = toc_items[items_per_column*2:]

    # insert into placeholders
    insert_toc_into_textframe(left_tf, left, template_para=left_para)
    insert_toc_into_textframe(right_tf, right, template_para=right_para)

    return leftovers


def handle_toc_multi_slides(prs, toc_items, items_per_column=30):
    """Distribute TOC across multiple slides until all items are placed."""
    # find untouched template
    toc_template = None
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame and "{{Table_Contents_Left}}" in shape.text:
                toc_template = slide
                break
        if toc_template:
            break
    if not toc_template:
        return

    current_items = toc_items
    slides_made = 0

    while current_items:
        # always duplicate the original template
        slide = duplicate_slide(prs, toc_template)
        leftovers = replace_toc_in_slide(slide, current_items, items_per_column)

        if leftovers == current_items:  # safeguard
            print("âš ï¸ Stopped early, items not consumed:", leftovers[:5])
            break

        current_items = leftovers
        slides_made += 1
        print(f"TOC Slide {slides_made}: {len(current_items)} items left")

    # finally remove the untouched template
    remove_slide(prs, toc_template)
    print(f"TOC generated across {slides_made} slides")

def remove_slide(prs, slide):
    """Safely remove a slide from a presentation."""
    slide_id = prs.slides._sldIdLst[prs.slides.index(slide)]
    prs.slides._sldIdLst.remove(slide_id)

# --------------- functions that modify PPT content ---------------

def insert_bullets_into_textframe(text_frame, items, template_para=None):
    """
    Clear the text_frame and insert one paragraph per item, prefixed by '•'.
    If template_para is provided, copy its first-run font + alignment to new paragraphs.
    """
    # choose template paragraph and run for formatting
    tmpl_para = template_para if template_para is not None else (text_frame.paragraphs[0] if text_frame.paragraphs else None)
    tmpl_run = None
    tmpl_align = None
    if tmpl_para is not None:
        tmpl_align = getattr(tmpl_para, "alignment", None)
        if tmpl_para.runs:
            tmpl_run = tmpl_para.runs[0]

    # clear the frame (reset to single empty paragraph)
    text_frame.text = ""

    for i, item in enumerate(items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        # use bullet character so we don't depend on PPT list styles
        p.text = "\u2022 " + item  # bullet + space
        if tmpl_align is not None:
            try:
                p.alignment = tmpl_align
            except Exception:
                pass
        # copy font formatting from template run (if any)
        if tmpl_run and p.runs:
            src_font = tmpl_run.font
            dst_run = p.runs[0]
            safe_copy_font(src_font, dst_run.font)

def distribute_items_across_cells(row, start_cell_index, items):
    """
    Distribute items across adjacent cells in row starting at start_cell_index.
    Logic:
     - if items <= 3 or only 1 available cell -> put all in start cell
     - if items > 3 and available_cells >= 2 -> split into 2 columns
     - if items > 6 and available_cells >= 3 -> split into 3 columns
    """
    available = len(row.cells) - start_cell_index
    n = len(items)
    # choose split count
    if n <= 3 or available == 1:
        splits = 1
    elif n > 6 and available >= 3:
        splits = 3
    else:
        splits = 2
    # compute even-ish distribution
    base = n // splits
    extra = n % splits
    chunks = []
    start = 0
    for i in range(splits):
        size = base + (1 if i < extra else 0)
        chunks.append(items[start:start+size])
        start += size
    # return list of (cell_index, chunk_items)
    result = []
    for i, chunk in enumerate(chunks):
        cell_index = start_cell_index + i
        if cell_index < len(row.cells):
            result.append((cell_index, chunk))
    return result

def set_text_with_format(shape, text, template_para):
    """
    Replace text inside shape but preserve formatting from template_para.
    """
    tf = shape.text_frame
    tf.clear()  # clear but keep text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    if template_para.runs:
        safe_copy_font(template_para.runs[0].font, run.font)

def set_text_with_placeholder_format(shape, text):
    """Replace shape text but preserve formatting of the placeholder's first run."""
    tf = shape.text_frame
    if not tf.paragraphs or not tf.paragraphs[0].runs:
        tf.text = text
        return

    # Copy font from first run
    src_font = tf.paragraphs[0].runs[0].font

    # Clear and rebuild with same formatting
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    safe_copy_font(src_font, run.font)

def clone_shape_style(src_shape, dst_shape):
    """Clone fill, line, alignment and margin from src_shape to dst_shape."""
    try:
        if src_shape.fill.type is not None:
            dst_shape.fill.fore_color.rgb = src_shape.fill.fore_color.rgb
    except:
        pass
    try:
        if src_shape.line and src_shape.line.color:
            dst_shape.line.color.rgb = src_shape.line.color.rgb
    except:
        pass

    if src_shape.has_text_frame and dst_shape.has_text_frame:
        dst_tf, src_tf = dst_shape.text_frame, src_shape.text_frame
        dst_tf.margin_left = src_tf.margin_left
        dst_tf.margin_right = src_tf.margin_right
        dst_tf.margin_top = src_tf.margin_top
        dst_tf.margin_bottom = src_tf.margin_bottom
        dst_tf.word_wrap = src_tf.word_wrap
        if src_tf.paragraphs and src_tf.paragraphs[0].alignment:
            dst_tf.paragraphs[0].alignment = src_tf.paragraphs[0].alignment

def get_placeholder_color(shape):
    """Fetch font color (RGB or theme) from the placeholder shape."""
    if shape.has_text_frame:
        tf = shape.text_frame
        if tf.paragraphs and tf.paragraphs[0].runs:
            font = tf.paragraphs[0].runs[0].font
            if font.color:
                if getattr(font.color, "rgb", None):
                    return ("rgb", font.color.rgb)
                elif getattr(font.color, "theme_color", None):
                    return ("theme", font.color.theme_color)
    return None

def get_placeholder_formatting(shape_or_cell):
    """Enhanced formatting extraction that captures more details"""
    formatting = {
        'font_name': None,
        'font_size': None,
        'font_bold': None,
        'font_italic': None,
        'font_underline': None,
        'font_color_rgb': None,
        'font_color_theme': None,
        'alignment': None,
        'line_spacing': None,
        'space_before': None,
        'space_after': None
    }
    
    # Check if the input is a table cell (_Cell) or a shape
    from pptx.table import _Cell
    is_cell = isinstance(shape_or_cell, _Cell)
    
    # Get the text frame
    tf = shape_or_cell.text_frame if is_cell else (
        shape_or_cell.text_frame if hasattr(shape_or_cell, 'has_text_frame') and shape_or_cell.has_text_frame else None
    )
    
    if tf and tf.paragraphs and tf.paragraphs[0].runs:
        para = tf.paragraphs[0]
        run = para.runs[0]
        font = run.font
        
        # Extract font properties with better null checking
        formatting['font_name'] = getattr(font, 'name', None)
        formatting['font_size'] = getattr(font, 'size', None)
        formatting['font_bold'] = getattr(font, 'bold', None)
        formatting['font_italic'] = getattr(font, 'italic', None)
        formatting['font_underline'] = getattr(font, 'underline', None)
        
        # Extract color more safely
        if hasattr(font, 'color') and font.color:
            if hasattr(font.color, 'rgb') and font.color.rgb:
                formatting['font_color_rgb'] = font.color.rgb
            elif hasattr(font.color, 'theme_color') and font.color.theme_color is not None:
                from pptx.enum.dml import MSO_THEME_COLOR_INDEX
                if font.color.theme_color != MSO_THEME_COLOR_INDEX.NOT_THEME_COLOR:
                    formatting['font_color_theme'] = font.color.theme_color
        
        # Extract paragraph properties
        formatting['alignment'] = getattr(para, 'alignment', None)
        formatting['line_spacing'] = getattr(para, 'line_spacing', None)
        formatting['space_before'] = getattr(para, 'space_before', None)
        formatting['space_after'] = getattr(para, 'space_after', None)
    
    return formatting

def apply_formatting_to_paragraph(paragraph, formatting):
    """Apply paragraph-level formatting"""
    if not formatting:
        return
        
    try:
        if formatting.get('alignment') is not None:
            paragraph.alignment = formatting['alignment']
        if formatting.get('line_spacing') is not None:
            paragraph.line_spacing = formatting['line_spacing']
        if formatting.get('space_before') is not None:
            paragraph.space_before = formatting['space_before']
        if formatting.get('space_after') is not None:
            paragraph.space_after = formatting['space_after']
    except Exception as e:
        print(f"Warning: Error applying paragraph formatting: {e}")

def apply_formatting_to_run(run, formatting):
    """Enhanced formatting application"""
    if not formatting:
        return
    
    font = run.font
    
    # Apply font properties only if they exist in formatting
    if formatting.get('font_name'):
        font.name = formatting['font_name']
    if formatting.get('font_size') is not None:
        font.size = formatting['font_size']
    if formatting.get('font_bold') is not None:
        font.bold = formatting['font_bold']
    if formatting.get('font_italic') is not None:
        font.italic = formatting['font_italic']
    if formatting.get('font_underline') is not None:
        font.underline = formatting['font_underline']
    
    # Apply color - RGB takes priority
    if formatting.get('font_color_rgb'):
        font.color.rgb = formatting['font_color_rgb']
    elif formatting.get('font_color_theme') is not None:
        font.color.theme_color = formatting['font_color_theme']

# NEW FUNCTION: Replace inline placeholders with comma-separated text
def replace_inline_placeholder_in_slide(slide, placeholder, items_or_text):
    """Inline replacement that preserves formatting by editing runs only."""
    inline_text = ", ".join(items_or_text) if isinstance(items_or_text, list) else str(items_or_text or "")
    replace_text_placeholders_in_slide(slide, placeholder, inline_text)

def get_table_cell_formatting(cell):
    """Extract comprehensive formatting from a table cell"""
    formatting = {
        'font_name': None,
        'font_size': None,
        'font_bold': None,
        'font_italic': None,
        'font_underline': None,
        'font_color_rgb': None,
        'font_color_theme': None,
        'alignment': None,
        'line_spacing': None,
        'space_before': None,
        'space_after': None,
        'margin_left': None,
        'margin_right': None,
        'margin_top': None,
        'margin_bottom': None
    }
    
    if not cell or not getattr(cell, "text_frame", None):
        return formatting
    
    tf = cell.text_frame
    
    # Extract text frame margin settings
    try:
        formatting['margin_left'] = tf.margin_left
        formatting['margin_right'] = tf.margin_right  
        formatting['margin_top'] = tf.margin_top
        formatting['margin_bottom'] = tf.margin_bottom
    except:
        pass
    
    # Extract formatting from first paragraph and run
    if tf.paragraphs and tf.paragraphs[0].runs:
        para = tf.paragraphs[0]
        run = para.runs[0]
        font = run.font
        
        # Font properties
        formatting['font_name'] = getattr(font, 'name', None)
        formatting['font_size'] = getattr(font, 'size', None)
        formatting['font_bold'] = getattr(font, 'bold', None) 
        formatting['font_italic'] = getattr(font, 'italic', None)
        formatting['font_underline'] = getattr(font, 'underline', None)
        
        # Color handling
        if hasattr(font, 'color') and font.color:
            if hasattr(font.color, 'rgb') and font.color.rgb:
                formatting['font_color_rgb'] = font.color.rgb
            elif hasattr(font.color, 'theme_color') and font.color.theme_color is not None:
                from pptx.enum.dml import MSO_THEME_COLOR_INDEX
                if font.color.theme_color != MSO_THEME_COLOR_INDEX.NOT_THEME_COLOR:
                    formatting['font_color_theme'] = font.color.theme_color
        
        # Paragraph properties
        formatting['alignment'] = getattr(para, 'alignment', None)
        formatting['line_spacing'] = getattr(para, 'line_spacing', None)
        formatting['space_before'] = getattr(para, 'space_before', None)
        formatting['space_after'] = getattr(para, 'space_after', None)
    
    return formatting

def apply_formatting_to_table_cell(cell, formatting, items, use_bullets=True):
    """Apply formatting to table cell and populate with items"""
    if not cell or not cell.text_frame or not formatting:
        return
    
    tf = cell.text_frame
    
    # Apply text frame margins if they were captured
    try:
        if formatting.get('margin_left') is not None:
            tf.margin_left = formatting['margin_left']
        if formatting.get('margin_right') is not None:
            tf.margin_right = formatting['margin_right']
        if formatting.get('margin_top') is not None:
            tf.margin_top = formatting['margin_top']
        if formatting.get('margin_bottom') is not None:
            tf.margin_bottom = formatting['margin_bottom']
    except:
        pass
    
    # Clear and rebuild content
    tf.clear()
    
    # Add each item as a separate paragraph (with or without bullet)
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # Create run with or without bullet based on use_bullets parameter
        run = p.add_run()
        if use_bullets:
            run.text = f"• {item}"
        else:
            run.text = item
        
        # Apply font formatting
        font = run.font
        if formatting.get('font_name'):
            font.name = formatting['font_name']
        if formatting.get('font_size') is not None:
            font.size = formatting['font_size']
        if formatting.get('font_bold') is not None:
            font.bold = formatting['font_bold']
        if formatting.get('font_italic') is not None:
            font.italic = formatting['font_italic']
        if formatting.get('font_underline') is not None:
            font.underline = formatting['font_underline']
        
        # Apply color
        if formatting.get('font_color_rgb'):
            font.color.rgb = formatting['font_color_rgb']
        elif formatting.get('font_color_theme') is not None:
            font.color.theme_color = formatting['font_color_theme']
        
        # Apply paragraph formatting
        if formatting.get('alignment') is not None:
            p.alignment = formatting['alignment']
        if formatting.get('line_spacing') is not None:
            p.line_spacing = formatting['line_spacing']
        if formatting.get('space_before') is not None:
            p.space_before = formatting['space_before']
        if formatting.get('space_after') is not None:
            p.space_after = formatting['space_after']

def replace_list_placeholder_in_table_with_expansion_enhanced(slide, placeholder, items, excel_path):
    """
    Enhanced version that handles both inline and row-expansion behavior with additional data columns
    """
    # Check if this is an expansion placeholder
    is_expand = placeholder.endswith("_EXPAND}}")
    
    for shape in slide.shapes:
        if not shape.has_table:
            continue
            
        table = shape.table
        placeholder_found = False
        
        # Find the placeholder in the table
        for row_idx, row in enumerate(table.rows):
            for col_idx, cell in enumerate(row.cells):
                if not cell.text_frame:
                    continue
                    
                # Check if this cell contains the placeholder
                full_text = ""
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text:
                            full_text += run.text
                
                if placeholder not in full_text:
                    continue
                
                placeholder_found = True
                print(f"Found table placeholder: {placeholder} with {len(items)} items (expand: {is_expand})")
                
                if is_expand and len(items) > 1:
                    # ENHANCED EXPANSION MODE: Create new rows with additional data columns
                    handle_table_row_expansion_enhanced(table, row_idx, col_idx, cell, items, placeholder, excel_path)
                else:
                    # INLINE MODE: Keep all items in one cell
                    handle_table_inline_replacement(cell, items, placeholder, use_bullets=True)
                
                return  # Found and processed, exit
        
        if placeholder_found:
            break

def find_header_row_volume_only(ws, search_year=2024):
    """
    Find the row index that contains the given year in the SALES VOLUME section only.
    Excludes percentage sections.
    """
    for i, row in enumerate(ws.iter_rows(min_row=1, max_row=20, values_only=True), start=1):
        values = [str(c) if c else "" for c in row]
        row_text = " ".join(values)
        
        # Skip percentage sections
        if "(%" in row_text or " %" in row_text or "Volume Share" in row_text:
            continue
            
        # Look for the year
        if any(str(search_year) in v for v in values):
            return i, values
    
    raise ValueError(f"Header row with {search_year} not found in volume section of {ws.title}")

def get_sheet_data_for_year(excel_path, sheet_name, year):
    """
    Extract data for a specific year from a sheet.
    Specifically looks for the SALES VOLUME section (not percentage section).
    Returns a dictionary {item_name: value}
    """
    wb = openpyxl.load_workbook(excel_path, data_only=True)
    if sheet_name not in wb.sheetnames:
        return {}
    
    ws = wb[sheet_name]
    
    # First, try to find the sales volume section (not the percentage section)
    # Look for a header that contains the year but NOT "(%)""
    header_row_idx = None
    year_col_idx = None
    
    for row_idx in range(1, 20):  # Search first 20 rows
        row = list(ws.iter_rows(min_row=row_idx, max_row=row_idx, values_only=True))[0]
        
        # Skip rows that contain "%" - we want volume data, not percentage data
        row_text = " ".join(str(cell) if cell else "" for cell in row)
        if "(%" in row_text or " %" in row_text:
            continue
            
        # Look for the year in this row
        for col_idx, cell in enumerate(row):
            if cell and str(year) in str(cell):
                header_row_idx = row_idx
                year_col_idx = col_idx
                break
        
        if header_row_idx and year_col_idx is not None:
            break
    
    if header_row_idx is None or year_col_idx is None:
        print(f"Could not find year {year} in sales volume section of {sheet_name}")
        return {}
    
    print(f"Found {year} data in {sheet_name} at row {header_row_idx}, column {year_col_idx}")
    
    # Extract data from rows after the header
    data = {}
    for row_idx in range(header_row_idx + 1, min(header_row_idx + 10, ws.max_row + 1)):
        row = list(ws.iter_rows(min_row=row_idx, max_row=row_idx, values_only=True))[0]
        
        # Stop if we hit an empty row or a row starting with "Total"
        if not row[0] or "Total" in str(row[0]):
            break
            
        # Stop if we hit another section (like percentage section)
        row_text = " ".join(str(cell) if cell else "" for cell in row)
        if "(%" in row_text or "Volume Share" in row_text or "Market Breakup" in row_text:
            break
        
        item_name = str(row[0]).strip()
        try:
            value = float(row[year_col_idx]) if row[year_col_idx] is not None else 0
            data[item_name] = value
            print(f"  {item_name}: {value}")
        except (ValueError, TypeError, IndexError):
            continue
    
    return data

def get_cagr_for_item(excel_path, sheet_name, item_name, start_year=2025, end_year=2033):
    """
    Calculate or extract CAGR for a specific item between two years.
    First tries to find a CAGR column, then calculates if data is available.
    """
    wb = openpyxl.load_workbook(excel_path, data_only=True)
    if sheet_name not in wb.sheetnames:
        return ""
    
    ws = wb[sheet_name]
    
    # Try to find existing CAGR column
    for row in ws.iter_rows(min_row=1, max_row=5, values_only=True):
        for col_idx, cell in enumerate(row):
            if cell and "CAGR" in str(cell) and f"{start_year}-{end_year}" in str(cell):
                # Found CAGR column, find the item's CAGR
                for data_row in ws.iter_rows(min_row=6, values_only=True):
                    if data_row[0] and str(data_row[0]).strip() == item_name:
                        cagr_val = data_row[col_idx] if col_idx < len(data_row) else None
                        if cagr_val:
                            try:
                                return f"{float(cagr_val)*100:.1f}%" if abs(float(cagr_val)) < 1 else f"{float(cagr_val):.1f}%"
                            except:
                                return str(cagr_val)
    
    # If no CAGR column found, try to calculate from start and end year data
    try:
        start_data = get_sheet_data_for_year(excel_path, sheet_name, start_year)
        end_data = get_sheet_data_for_year(excel_path, sheet_name, end_year)
        
        if item_name in start_data and item_name in end_data:
            start_val = start_data[item_name]
            end_val = end_data[item_name]
            
            if start_val > 0 and end_val > 0:
                years = end_year - start_year
                cagr = ((end_val / start_val) ** (1/years) - 1) * 100
                return f"{cagr:.1f}%"
    except:
        pass
    
    return ""

def handle_table_row_expansion_enhanced(table, template_row_idx, col_idx, template_cell, items, placeholder, excel_path):
    """
    Enhanced version that adapts to table column count:
    - 2 columns: Item Name, 2024 Value
    - 3+ columns: Item Name, Unit, 2024 Value, 2033 Value, CAGR
    """
    
    # Extract formatting from the template cell
    template_formatting = get_table_cell_formatting(template_cell)
    
    # Get the number of columns in the table
    num_columns = len(table.rows[template_row_idx].cells)
    available_columns = num_columns - col_idx
    
    # Get unit from Summary sheet
    wb = openpyxl.load_workbook(excel_path, data_only=True)
    unit = ""
    if "Summary" in wb.sheetnames:
        summary_ws = wb["Summary"]
        for row in summary_ws.iter_rows(values_only=True):
            if row[0] and str(row[0]).strip().lower() == "unit":
                unit = str(row[1]) if row[1] else ""
                break
    
    # Determine which sheet to get data from based on placeholder
    sheet_name = ""
    if "By_Type" in placeholder:
        sheet_name = "By_Type"
    elif "By_Application" in placeholder:
        sheet_name = "By_Application"
    elif "By_EndUser" in placeholder:
        sheet_name = "By_EndUser"
    elif "By_Region" in placeholder:
        sheet_name = "By_Region"
    
    # Get data for 2024 and 2033
    data_2024 = get_sheet_data_for_year(excel_path, sheet_name, 2024) if sheet_name else {}
    data_2033 = get_sheet_data_for_year(excel_path, sheet_name, 2033) if sheet_name else {}
    
    # Clear the template cell and put first item with appropriate data based on column count
    if items:
        first_item = items[0]
        
        # Get values for first item
        val_2024 = data_2024.get(first_item, 0)
        val_2033 = data_2033.get(first_item, 0)
        cagr = get_cagr_for_item(excel_path, sheet_name, first_item) if sheet_name else ""
        
        # Adapt row data based on available columns
        if available_columns == 2:
            # 2 columns: Item Name, 2024 Value
            row_data = [
                first_item,
                f"{val_2024:,.1f}" if val_2024 else ""
            ]
        elif available_columns == 3:
            # 3 columns: Item Name, Unit, 2024 Value
            row_data = [
                first_item,
                unit,
                f"{val_2024:,.1f}" if val_2024 else ""
            ]
        elif available_columns == 4:
            # 4 columns: Item Name, Unit, 2024 Value, 2033 Value
            row_data = [
                first_item,
                unit,
                f"{val_2024:,.1f}" if val_2024 else "",
                f"{val_2033:,.1f}" if val_2033 else ""
            ]
        else:
            # 5+ columns: Item Name, Unit, 2024 Value, 2033 Value, CAGR
            row_data = [
                first_item,
                unit,
                f"{val_2024:,.1f}" if val_2024 else "",
                f"{val_2033:,.1f}" if val_2033 else "",
                cagr
            ]
        
        fill_table_row_with_data(table.rows[template_row_idx], col_idx, row_data, template_formatting)
    
    # Add new rows for remaining items (items[1:])
    if len(items) > 1:
        for i, item in enumerate(items[1:], 1):
            try:
                # Clone the template row
                new_row = clone_table_row(table, template_row_idx)
                
                # Get values for this item
                val_2024 = data_2024.get(item, 0)
                val_2033 = data_2033.get(item, 0)
                cagr = get_cagr_for_item(excel_path, sheet_name, item) if sheet_name else ""
                
                # Adapt row data based on available columns
                if available_columns == 2:
                    # 2 columns: Item Name, 2024 Value
                    row_data = [
                        item,
                        f"{val_2024:,.1f}" if val_2024 else ""
                    ]
                elif available_columns == 3:
                    # 3 columns: Item Name, Unit, 2024 Value
                    row_data = [
                        item,
                        unit,
                        f"{val_2024:,.1f}" if val_2024 else ""
                    ]
                elif available_columns == 4:
                    # 4 columns: Item Name, Unit, 2024 Value, 2033 Value
                    row_data = [
                        item,
                        unit,
                        f"{val_2024:,.1f}" if val_2024 else "",
                        f"{val_2033:,.1f}" if val_2033 else ""
                    ]
                else:
                    # 5+ columns: Item Name, Unit, 2024 Value, 2033 Value, CAGR
                    row_data = [
                        item,
                        unit,
                        f"{val_2024:,.1f}" if val_2024 else "",
                        f"{val_2033:,.1f}" if val_2033 else "",
                        cagr
                    ]
                
                fill_table_row_with_data(new_row, col_idx, row_data, template_formatting)
                        
            except Exception as e:
                print(f"Error creating enhanced row {i}: {e}")
                break

def fill_table_row_with_data(row, start_col_idx, data_list, formatting):
    """
    Fill a table row starting from start_col_idx with data from data_list
    """
    for i, data in enumerate(data_list):
        col_idx = start_col_idx + i
        if col_idx < len(row.cells):
            cell = row.cells[col_idx]
            if cell.text_frame:
                cell.text_frame.clear()
                p = cell.text_frame.paragraphs[0]
                run = p.add_run()
                run.text = str(data)
                
                # Apply formatting
                if formatting:
                    apply_formatting_to_table_cell_content_single(run, p, formatting)

def handle_table_inline_replacement(cell, items, placeholder, use_bullets=True):
    """
    Default "inline" behavior for {{Key_List}}:
    • Renders all items as bullets in the same cell
    • Does NOT create new rows
    """
    if not cell or not cell.text_frame:
        return

    # Extract formatting from the template cell
    template_formatting = get_table_cell_formatting(cell)

    # Clear existing content
    tf = cell.text_frame
    tf.clear()

    # Add each item as a paragraph (with or without bullets)
    for item in items:
        p = tf.add_paragraph()
        p.level = 0
        run = p.add_run()
        if use_bullets:
            run.text = f"• {item}"
        else:
            run.text = item
        # Apply your stored formatting
        apply_formatting_to_table_cell_content_single(run, p, template_formatting)

def clone_table_row(table, template_row_idx):
    """Clone a table row and add it to the table"""
    from copy import deepcopy
    
    # Get the template row XML
    template_row_xml = table._tbl.tr_lst[template_row_idx]
    
    # Create a deep copy
    new_row_xml = deepcopy(template_row_xml)
    
    # Append to table
    table._tbl.append(new_row_xml)
    
    # Return the new row object (last row in table)
    return table.rows[len(table.rows) - 1]

def clear_cell_content(cell):
    """Clear content from a table cell while preserving structure"""
    if cell and cell.text_frame:
        cell.text_frame.clear()

def replace_cell_content_with_formatting(cell, items, formatting, use_bullets=True):
    """Replace cell content with formatted items"""
    if not cell or not cell.text_frame:
        return
    
    tf = cell.text_frame
    tf.clear()
    
    # Add each item as a paragraph (with or without bullet points)
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        run = p.add_run()
        if use_bullets:
            run.text = f"• {item}"
        else:
            run.text = item
        
        # Apply formatting from template
        apply_formatting_to_table_cell_content_single(run, p, formatting)

def apply_formatting_to_table_cell_content_single(run, paragraph, formatting):
    """Apply formatting to a single run and paragraph"""
    if not formatting:
        return
    
    # Apply font formatting to run
    font = run.font
    if formatting.get('font_name'):
        font.name = formatting['font_name']
    if formatting.get('font_size') is not None:
        font.size = formatting['font_size']
    if formatting.get('font_bold') is not None:
        font.bold = formatting['font_bold']
    if formatting.get('font_italic') is not None:
        font.italic = formatting['font_italic']
    
    # Apply color
    if formatting.get('font_color_rgb'):
        font.color.rgb = formatting['font_color_rgb']
    elif formatting.get('font_color_theme') is not None:
        font.color.theme_color = formatting['font_color_theme']
    
    # Apply paragraph formatting
    if formatting.get('alignment') is not None:
        paragraph.alignment = formatting['alignment']

def apply_formatting_to_table_cell_content(cell, formatting):
    """Apply formatting to entire cell content"""
    if not cell or not cell.text_frame or not formatting:
        return
    
    tf = cell.text_frame
    
    # Apply margins if available
    if formatting.get('margin_left') is not None:
        tf.margin_left = formatting['margin_left']
    if formatting.get('margin_right') is not None:
        tf.margin_right = formatting['margin_right']
    if formatting.get('margin_top') is not None:
        tf.margin_top = formatting['margin_top']
    if formatting.get('margin_bottom') is not None:
        tf.margin_bottom = formatting['margin_bottom']

# Modified main processing function
def process_table_placeholders_with_expansion_enhanced(slide, list_placeholders, excel_path):
    """
    Enhanced version that processes both regular and expansion table placeholders with additional data
    """
    for key, items in list_placeholders.items():
        # Regular placeholder (inline behavior)
        regular_placeholder = "{{" + key + "}}"
        replace_list_placeholder_in_table_with_expansion_enhanced(slide, regular_placeholder, items, excel_path)
        
        # Expansion placeholder (row creation behavior with enhanced data)
        expand_placeholder = "{{" + key + "_EXPAND}}"
        replace_list_placeholder_in_table_with_expansion_enhanced(slide, expand_placeholder, items, excel_path)

# UPDATED FUNCTION: Replace list placeholder by rebuilding text_frame
def replace_list_placeholder_in_slide(slide, placeholder, items):
    # Find the shape containing the placeholder
    targets = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        full = "".join((run.text or "") for p in shape.text_frame.paragraphs for run in p.runs)
        if placeholder in full:
            targets.append(shape)
    if not targets:
        return  # Placeholder not found

    # --- Case A: multiple shapes â†’ segmentation boxes ---
    if len(targets) > 1:
        tpl_fmt = get_placeholder_formatting(targets[0])
        for i, shape in enumerate(targets):
            tf = shape.text_frame
            tf.clear()
            if i < len(items):
                p = tf.paragraphs[0]
                r = p.add_run()
                r.text = str(items[i])
                apply_formatting_to_run(r, tpl_fmt)
                # Apply paragraph-level formatting too
                if tpl_fmt.get("alignment") is not None:
                    p.alignment = tpl_fmt["alignment"]
                # typical white text over colored boxes for segmentation
                try: 
                    r.font.color.rgb = RGBColor(255, 255, 255)
                except: 
                    pass
        # add extra boxes if more items than shapes
        if len(items) > len(targets):
            template = targets[-1]
            try:
                gap = (targets[1].top - targets[0].top) if len(targets) >= 2 else template.height + Emu(180000)
            except:
                gap = template.height + Emu(180000)
            for j, item in enumerate(items[len(targets):], start=1):
                try:
                    new_el = deepcopy(template.element)
                    slide.shapes._spTree.insert_element_before(new_el, "p:extLst")
                    new_shape = slide.shapes[-1]
                    new_shape.left = template.left
                    new_shape.top = template.top + gap * j
                    new_shape.width, new_shape.height = template.width, template.height
                    tf = new_shape.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    r = p.add_run()
                    r.text = str(item)
                    apply_formatting_to_run(r, tpl_fmt)
                    if tpl_fmt.get("alignment") is not None:
                        p.alignment = tpl_fmt["alignment"]
                    try: 
                        r.font.color.rgb = RGBColor(255, 255, 255)
                    except: 
                        pass
                except Exception as e:
                    print("Error cloning segmentation box:", e)
                    break
        return

    # --- Case B: single shape â†’ preserve bullet/indent + font exactly ---
    shape = targets[0]
    tf = shape.text_frame

    # find the paragraph that contains the placeholder
    target_para = None
    original_formatting = None
    for para in tf.paragraphs:
        para_text = "".join((run.text or "") for run in para.runs)
        if placeholder in para_text:
            target_para = para
            # IMPORTANT: Extract formatting from the placeholder paragraph itself
            if para.runs:
                original_formatting = get_run_formatting(para.runs[0])
            break
    
    if target_para is None:
        return

    # Split the text around the placeholder
    full_para_text = "".join(r.text or "" for r in target_para.runs)
    if placeholder not in full_para_text:
        return
        
    before_txt, after_txt = full_para_text.split(placeholder, 1)

    # Get paragraph-level formatting from the target paragraph
    para_formatting = {
        'alignment': getattr(target_para, 'alignment', None),
        'line_spacing': getattr(target_para, 'line_spacing', None),
        'space_before': getattr(target_para, 'space_before', None),
        'space_after': getattr(target_para, 'space_after', None),
        'level': getattr(target_para, 'level', None)
    }

    # Clear runs in the target paragraph but preserve paragraph structure
    for r in list(target_para.runs):
        r.text = ""
    
    # Add "before" text if exists
    if before_txt.strip():
        run = target_para.runs[0] if target_para.runs else target_para.add_run()
        run.text = before_txt
        if original_formatting:
            apply_run_formatting(run, original_formatting)

    # Insert each list item as a new paragraph with proper formatting
    for idx, item in enumerate(items):
        p = tf.add_paragraph()
        
        # Position the new paragraph after the target paragraph
        try:
            body = tf._txBody
            body.remove(p._p)
            body.insert(body.index(target_para._p) + 1 + idx, p._p)
        except Exception:
            pass
        
        # Apply paragraph-level formatting
        if para_formatting.get('alignment') is not None:
            p.alignment = para_formatting['alignment']
        if para_formatting.get('line_spacing') is not None:
            p.line_spacing = para_formatting['line_spacing']
        if para_formatting.get('space_before') is not None:
            p.space_before = para_formatting['space_before']
        if para_formatting.get('space_after') is not None:
            p.space_after = para_formatting['space_after']
        if para_formatting.get('level') is not None:
            p.level = para_formatting['level']
        
        # Create the run with bullet and text
        r = p.add_run()
        r.text = f"• {item}"
        
        # Apply font formatting from original placeholder
        if original_formatting:
            apply_run_formatting(r, original_formatting)

    # Add "after" text if exists
    if after_txt.strip():
        p = tf.add_paragraph()
        try:
            body = tf._txBody
            body.remove(p._p)
            body.insert(body.index(target_para._p) + 1 + len(items), p._p)
        except Exception:
            pass
        
        # Apply same paragraph formatting
        if para_formatting.get('alignment') is not None:
            p.alignment = para_formatting['alignment']
        if para_formatting.get('line_spacing') is not None:
            p.line_spacing = para_formatting['line_spacing']
        if para_formatting.get('space_before') is not None:
            p.space_before = para_formatting['space_before']
        if para_formatting.get('space_after') is not None:
            p.space_after = para_formatting['space_after']
        if para_formatting.get('level') is not None:
            p.level = para_formatting['level']
            
        r = p.add_run()
        r.text = after_txt
        if original_formatting:
            apply_run_formatting(r, original_formatting)

def get_paragraph_formatting(paragraph):
    """Extract formatting from a paragraph"""
    formatting = {
        'alignment': getattr(paragraph, 'alignment', None),
        'line_spacing': getattr(paragraph, 'line_spacing', None),
        'space_before': getattr(paragraph, 'space_before', None),
        'space_after': getattr(paragraph, 'space_after', None),
        'font_formatting': None
    }
    
    if paragraph.runs:
        font = paragraph.runs[0].font
        formatting['font_formatting'] = {
            'font_name': getattr(font, 'name', None),
            'font_size': getattr(font, 'size', None),
            'font_bold': getattr(font, 'bold', None),
            'font_italic': getattr(font, 'italic', None),
            'font_underline': getattr(font, 'underline', None),
        }
        
        # Handle color
        if hasattr(font, 'color') and font.color:
            if hasattr(font.color, 'rgb') and font.color.rgb:
                formatting['font_formatting']['font_color_rgb'] = font.color.rgb
            elif hasattr(font.color, 'theme_color') and font.color.theme_color is not None:
                formatting['font_formatting']['font_color_theme'] = font.color.theme_color
    
    return formatting

def apply_paragraph_formatting(paragraph, formatting):
    """Apply stored paragraph formatting"""
    if not formatting:
        return
        
    if formatting.get('alignment') is not None:
        paragraph.alignment = formatting['alignment']
    if formatting.get('line_spacing') is not None:
        paragraph.line_spacing = formatting['line_spacing']
    if formatting.get('space_before') is not None:
        paragraph.space_before = formatting['space_before']
    if formatting.get('space_after') is not None:
        paragraph.space_after = formatting['space_after']
    
    # Apply font formatting to runs
    font_fmt = formatting.get('font_formatting')
    if font_fmt and paragraph.runs:
        for run in paragraph.runs:
            font = run.font
            if font_fmt.get('font_name'):
                font.name = font_fmt['font_name']
            if font_fmt.get('font_size') is not None:
                font.size = font_fmt['font_size']
            if font_fmt.get('font_bold') is not None:
                font.bold = font_fmt['font_bold']
            if font_fmt.get('font_italic') is not None:
                font.italic = font_fmt['font_italic']
            if font_fmt.get('font_underline') is not None:
                font.underline = font_fmt['font_underline']
            
            if font_fmt.get('font_color_rgb'):
                font.color.rgb = font_fmt['font_color_rgb']
            elif font_fmt.get('font_color_theme') is not None:
                font.color.theme_color = font_fmt['font_color_theme']

def replace_text_placeholders_in_slide(slide, placeholder, replacement):
    """
    Enhanced version that better preserves formatting when replacing placeholders.
    Works by preserving the formatting of the run that contains the start of the placeholder.
    """
    if not placeholder:
        return

    def replace_in_paragraph_runs(paragraph):
        runs = paragraph.runs
        if not runs:
            return False

        # Build concatenated string and track run positions
        run_texts = [r.text or "" for r in runs]
        full = "".join(run_texts)
        
        if placeholder not in full:
            return False

        replaced = False
        
        # Process each occurrence of the placeholder
        while placeholder in full:
            start = full.index(placeholder)
            end = start + len(placeholder)

            # Find which run contains the start of the placeholder
            pos = 0
            start_run_idx = -1
            start_run_formatting = None
            
            for i, (run, text) in enumerate(zip(runs, run_texts)):
                if pos <= start < pos + len(text):
                    start_run_idx = i
                    start_run_formatting = get_run_formatting(run)
                    break
                pos += len(text)

            # Build run position metadata
            pos = 0
            runs_meta = []
            for i, (r, t) in enumerate(zip(runs, run_texts)):
                runs_meta.append({
                    "run": r, 
                    "text": t, 
                    "start": pos, 
                    "end": pos + len(t),
                    "index": i
                })
                pos += len(t)

            # Process each run that intersects with the placeholder
            new_runs = []
            current_pos = 0
            for meta in runs_meta:
                r = meta["run"]
                t = meta["text"]
                s = meta["start"]
                e = meta["end"]
                idx = meta["index"]

                # Skip runs that don't intersect with placeholder
                if e <= start or s >= end:
                    new_runs.append((r, t))
                    continue

                # Calculate intersection within this run
                run_rel_start = max(0, start - s)
                run_rel_end = min(len(t), end - s)
                
                before = t[:run_rel_start]
                after = t[run_rel_end:]

                # Determine new text for this run
                if s <= start < e:  # This run contains the start of placeholder
                    new_text = before + replacement
                    if end <= e:  # Placeholder also ends in this run
                        new_text += after
                    
                    # Update the run text and preserve formatting
                    r.text = new_text
                    if start_run_formatting and idx == start_run_idx:
                        apply_run_formatting(r, start_run_formatting)
                    new_runs.append((r, new_text))
                
                elif s < end <= e:  # This run contains the end of placeholder
                    r.text = after
                    new_runs.append((r, after))
                else:  # This run is completely inside the placeholder
                    r.text = ""
                    new_runs.append((r, ""))

            # Rebuild runs to ensure formatting is maintained across the replacement
            for i in reversed(range(len(paragraph.runs))):
                paragraph._element.remove(paragraph.runs[i]._r)

            for run, text in new_runs:
                if text:
                    new_run = paragraph.add_run()
                    new_run.text = text
                    apply_run_formatting(new_run, get_run_formatting(run))

            # Rebuild for next iteration
            run_texts = [r.text or "" for r in paragraph.runs]
            full = "".join(run_texts)
            replaced = True

        return replaced

    # Process all shapes in the slide
    for shape in slide.shapes:
        # Handle text frames
        if getattr(shape, "has_text_frame", False):
            tf = shape.text_frame
            for para in tf.paragraphs:
                replace_in_paragraph_runs(para)

        # Handle tables
        if getattr(shape, "has_table", False):
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    if not getattr(cell, "text_frame", None):
                        continue
                    for para in cell.text_frame.paragraphs:
                        replace_in_paragraph_runs(para)

def add_row_to_table(table, template_row_idx):
    """Clone a row in the table at the end, using template_row_idx as format."""
    from copy import deepcopy
    row_xml = table._tbl.tr_lst[template_row_idx]
    new_row = deepcopy(row_xml)
    table._tbl.append(new_row)
    # Table.rows does not support negative indexing; use explicit last index
    return table.rows[len(table.rows) - 1]

def _cell_contains_placeholder(cell, placeholder):
    """Return True if any runs in any paragraph of this cell contain placeholder text."""
    if not cell or not getattr(cell, "text_frame", None):
        return False
    full = ""
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            if run.text:
                full += run.text
    return placeholder in full

def delete_extra_rows(tbl, start_index, keep_rows):
    """
    Delete extra rows from PPTX table by editing the underlying XML.
    Works without leaving blanks.
    """
    # All rows as XML
    tr_elements = list(tbl._tbl.findall('.//a:tr', namespaces={'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}))

    target_end = start_index + keep_rows
    if len(tr_elements) <= target_end:
        return

    # remove rows after target_end
    for tr in tr_elements[target_end:]:
        tbl._tbl.remove(tr)

def normalize_products(raw):
    """
    Return a list of product strings.
    Accepts: list, JSON list string, comma/semicolon/newline-separated string.
    """
    if not raw:
        return []
    if isinstance(raw, list):
        return [str(x).strip() for x in raw if str(x).strip()]

    s = str(raw).strip()
    # try JSON parse
    try:
        parsed = json.loads(s)
        if isinstance(parsed, list):
            return [str(x).strip() for x in parsed if str(x).strip()]
    except Exception:
        pass

    # split on common delimiters (newline, bullet, comma, semicolon)
    parts = re.split(r'[\n\u2022•;]+|,\s*', s)
    cleaned = [p.strip().strip('"').strip("'") for p in parts if p.strip()]
    return cleaned

# --- helper: quick Wikipedia-based founding year lookup (fallback) ---
def fetch_founding_from_wikipedia(company_name, timeout=8):
    """
    Try to find a founding year from the company's Wikipedia page.
    Returns a string year (e.g. "1897") or "" if not found.
    """
    try:
        API = "https://en.wikipedia.org/w/api.php"
        # 1) search for the page
        params = {"action": "query", "list": "search", "srsearch": company_name, "format": "json", "srlimit": 1}
        r = requests.get(API, params=params, timeout=timeout)
        r.raise_for_status()
        data = r.json()
        hits = data.get("query", {}).get("search", [])
        if not hits:
            return ""
        title = hits[0]["title"]

        # 2) fetch plaintext extract
        params2 = {"action": "query", "prop": "extracts", "explaintext": 1, "titles": title, "format": "json", "redirects": 1}
        r2 = requests.get(API, params=params2, timeout=timeout)
        r2.raise_for_status()
        pages = r2.json().get("query", {}).get("pages", {})
        if not pages:
            return ""

        # get extract text
        page = next(iter(pages.values()))
        extract = page.get("extract", "")

        # look for 'Founded', 'Established', 'founded in' patterns
        m = re.search(r'(?:Founded|Founded in|founded in|Established|established|Founded:|Founded -)\D{0,30}(\d{4})', extract, re.I)
        if m:
            year = m.group(1)
            return year

        # fallback: first plausible 4-digit year in whole extract
        m2 = re.search(r'(\b(17|18|19|20)\d{2}\b)', extract)
        if m2:
            y = int(m2.group(1))
            if 1700 <= y <= datetime.datetime.now().year:
                return str(y)

    except Exception:
        return ""
    return ""

def fetch_company_details(company_name, use_ai=True, ai_timeout=8):
    details = {
        "founding_year": "",
        "headquarters": "",
        "website": "",
        "products_offered": []
    }

    if use_ai:
        try:
            prompt = f"""
            Provide very short structured details about the company "{company_name}".
            Return JSON with keys: founding_year, headquarters, website, products_offered.
            Example:
            {{
              "founding_year": "1897",
              "headquarters": "Tokyo, Japan",
              "website": "https://www.example.com",
              "products_offered": ["Chemicals", "Plastics"]
            }}
            Keep it concise and factual.
            """
            model = genai.GenerativeModel("gemini-1.5-flash")
            response = model.generate_content(prompt)

            text = getattr(response, "text", "").strip()

            # --- Clean Gemini output ---
            # remove code fences if present
            if text.startswith("```"):
                text = re.sub(r"^```[a-zA-Z]*\n", "", text)
                text = re.sub(r"\n```$", "", text)
                text = text.strip()

            # Try parsing JSON
            parsed = {}
            try:
                parsed = json.loads(text)
            except Exception:
                # fallback: try to extract {...} JSON substring
                m = re.search(r"\{.*\}", text, flags=re.S)
                if m:
                    try:
                        parsed = json.loads(m.group(0))
                    except:
                        pass

            # fallback if still not JSON
            if not isinstance(parsed, dict):
                parsed = {}

            # merge parsed fields
            details["founding_year"] = str(parsed.get("founding_year", "") or "")
            details["headquarters"] = str(parsed.get("headquarters", "") or "")
            details["website"] = str(parsed.get("website", "") or "")
            details["products_offered"] = parsed.get("products_offered", [])
        except Exception as e:
            # non-fatal - we will try fallbacks below
            print("⚠️ AI lookup failed or timed out:", e)

    # 2) Normalize products into a list
    details["products_offered"] = normalize_products(details.get("products_offered", ""))

    # 3) Validate founding_year — must be a 4-digit plausible year
    fy_raw = details.get("founding_year", "")
    fy_candidate = None
    if fy_raw:
        m = re.search(r'(\d{4})', str(fy_raw))
        if m:
            try:
                fy_candidate = int(m.group(1))
            except Exception:
                fy_candidate = None

    now_year = datetime.datetime.now().year
    if not fy_candidate or fy_candidate < 1700 or fy_candidate > now_year:
        # fallback to Wikipedia
        wiki_year = fetch_founding_from_wikipedia(company_name)
        if wiki_year:
            details["founding_year"] = wiki_year
        else:
            # if AI gave something non-plausible, blank it
            details["founding_year"] = "" if not fy_candidate else str(fy_candidate)
    else:
        details["founding_year"] = str(fy_candidate)

    # 4) final normalization: ensure products_offered is list of strings
    details["products_offered"] = normalize_products(details.get("products_offered", ""))

    return details

def distribute_company_names_across_template_slides(prs, placeholder, items, duplicate_if_needed=True):
    """
    Fill company details dynamically in table (using Gemini for details).
    Columns assumed as:
      col_idx = Company Name
      col_idx+1 = Founding Year
      col_idx+2 = Headquarters
      col_idx+3 = Website
      col_idx+4 = Products Offered
    """

    # 1) collect templates (in slide order)
    templates = []
    for s_idx, slide in enumerate(prs.slides):
        for shp in slide.shapes:
            if not getattr(shp, "has_table", False):
                continue
            tbl = shp.table
            for r_idx, row in enumerate(tbl.rows):
                for c_idx, cell in enumerate(row.cells):
                    if _cell_contains_placeholder(cell, placeholder):
                        formatting = get_placeholder_formatting(cell)
                        is_header = (r_idx == 0)
                        header_offset = r_idx + 1 if is_header else r_idx
                        capacity = max(0, len(tbl.rows) - header_offset)

                        # clear placeholder
                        cell.text = ""

                        templates.append({
                            "slide_idx": s_idx,
                            "slide": slide,
                            "shape": shp,
                            "col_idx": c_idx,
                            "header_offset": header_offset,
                            "capacity": capacity,
                            "formatting": formatting
                        })
                        break

    if not templates:
        print(f"âš ï¸ No template slides found containing {placeholder}. Nothing filled.")
        return

    total_capacity = sum(t["capacity"] for t in templates)
    print(f"Found {len(templates)} template slide(s), capacities:", [t["capacity"] for t in templates], "â†’ total {total_capacity}")

    def _fill_table_object(tbl_obj, col_idx, header_offset, chunk_items, formatting):
        """
        Fills a chunk of companies into table tbl_obj starting at header_offset row,
        writes company name, founding_year, hq, website, and BULLETED products in the last cell.
        """
        for i, company in enumerate(chunk_items):
            row_index = i + header_offset
            if row_index >= len(tbl_obj.rows):
                continue

            # Get details (AI + fallback)
            details = fetch_company_details(company, use_ai=True)

            # Values for first four columns
            values = [
                company,
                details.get("founding_year", ""),
                details.get("headquarters", ""),
                details.get("website", "")
            ]

            # Fill the first four columns (company, founding, hq, website)
            for j, val in enumerate(values):
                try:
                    cell_tf = tbl_obj.cell(row_index, col_idx + j).text_frame
                    cell_tf.clear()
                    p = cell_tf.paragraphs[0]
                    run = p.add_run()
                    run.text = str(val)
                    apply_formatting_to_paragraph(p, formatting)
                except Exception as e:
                    print(f"âš ï¸ Error filling cell ({row_index},{col_idx + j}): {e}")

            # Fill products column as distinct bulleted paragraphs
            try:
                prod_cell = tbl_obj.cell(row_index, col_idx + 4)
                prod_tf = prod_cell.text_frame
                prod_tf.clear()
                products = details.get("products_offered", [])
                if isinstance(products, str):
                    products = normalize_products(products)

                if not products:
                    # keep cell empty
                    pass
                else:
                    for k, prod in enumerate(products):
                        p = prod_tf.paragraphs[0] if k == 0 else prod_tf.add_paragraph()
                        run = p.add_run()
                        run.text = f"• {prod}"
                        apply_formatting_to_paragraph(p, formatting)
            except Exception as e:
                print(f"âš ï¸ Error filling products cell ({row_index},{col_idx + 4}): {e}")

        # After filling N items, delete the extra rows (keep rows = number items)
        delete_extra_rows(tbl_obj, header_offset, len(chunk_items))

    # 2) iterate templates and fill
    i = 0
    for t in templates:
        if i >= len(items):
            break
        cap = t["capacity"]
        if cap <= 0:
            continue
        chunk = items[i:i+cap]
        _fill_table_object(t["shape"].table, t["col_idx"], t["header_offset"], chunk, t["formatting"])
        print(f"Filled slide {t['slide_idx']} with {len(chunk)} companies")
        i += len(chunk)

    # 3) handle leftovers â†’ duplicate slides
    while i < len(items):
        if not duplicate_if_needed:
            break
        last_template = templates[-1]
        new_slide = duplicate_slide(prs, last_template["slide"])

        target_tbl = None
        for shp in new_slide.shapes:
            if getattr(shp, "has_table", False):
                target_tbl = shp.table
                break
        if not target_tbl:
            break

        cap = max(0, len(target_tbl.rows) - last_template["header_offset"])
        if cap <= 0:
            break
        chunk = items[i:i+cap]
        _fill_table_object(target_tbl, last_template["col_idx"], last_template["header_offset"], chunk, last_template["formatting"])
        print(f"Filled duplicated slide with {len(chunk)} companies")
        i += len(chunk)

    print(f"Done filling companies. Total filled: {min(len(items), i)} / {len(items)}")

def update_charts_in_slide(slide, volumes, unit, historical_years=[2019, 2020, 2021, 2022, 2023, 2024], forecast_years=[2025, 2026, 2027, 2028, 2029, 2030, 2031, 2032, 2033]):
    """
    Update bar chart series data from volumes dict.
    Assumes single series per chart; historical/forecast based on categories containing years.
    """
    for shape in slide.shapes:
        if not shape.has_chart:
            continue
        chart = shape.chart
        if chart.chart_type != XL_CHART_TYPE.COLUMN_CLUSTERED:
            continue
        
        # Get categories from the chart's plots - correct way to access categories
        categories = []
        try:
            if chart.plots and len(chart.plots) > 0:
                plot = chart.plots[0]
                if hasattr(plot, 'categories') and plot.categories:
                    categories = [str(cat) for cat in plot.categories]
        except AttributeError:
            # Fallback: try to get categories from series if available
            try:
                if chart.series and len(chart.series) > 0:
                    # Categories might be accessible through chart data
                    pass
            except:
                pass
        
        title_text = ""
        try:
            if chart.has_title and chart.chart_title:
                title_text = chart.chart_title.text if chart.chart_title.text else ""
        except:
            pass
        
        # Determine if this is historical or forecast chart
        is_historical = any(str(y) in ' '.join(categories + [title_text]) for y in historical_years)
        is_forecast = any(str(y) in ' '.join(categories + [title_text]) for y in forecast_years)
        
        if not (is_historical or is_forecast):
            continue
        
        years = historical_years if is_historical else forecast_years
        
        # Update existing series values
        try:
            if chart.series and len(chart.series) > 0:
                series = chart.series[0]
                # Update the values
                new_values = [volumes.get(y, 0) for y in years]
                
                # Try to update values directly
                try:
                    series.values = new_values
                    print(f"Updated {'historical' if is_historical else 'forecast'} chart series with {len(years)} data points")
                except Exception as e:
                    print(f"Warning: Could not update series values directly: {e}")
                    
                    # Fallback: replace entire chart data
                    try:
                        data = CategoryChartData()
                        data.categories = [str(y) for y in years]
                        data.add_series("Volume", tuple(new_values))
                        chart.replace_data(data)
                        print(f"Replaced {'historical' if is_historical else 'forecast'} chart data with {len(years)} data points")
                    except Exception as e2:
                        print(f"Warning: Could not replace chart data: {e2}")
            else:
                # No existing series, create new data
                data = CategoryChartData()
                data.categories = [str(y) for y in years]
                data.add_series("Volume", tuple(volumes.get(y, 0) for y in years))
                chart.replace_data(data)
                print(f"Created new {'historical' if is_historical else 'forecast'} chart with {len(years)} data points")
                
        except Exception as e:
            print(f"Error updating chart: {e}")
            continue

# --------------- main ---------------

def main(excel_file, ppt_template, output_ppt):
    prs = Presentation(ppt_template)

    kv = read_summary_keys(excel_file, "Summary")
    dynamic_kv, volumes = extract_dynamic_placeholders(excel_file, include_market_overview=True, include_overview_content=True)
    kv.update(dynamic_kv)
    kv["Subtitle"] = build_report_subtitle(excel_file)

    # Create inline versions of list placeholders
    inline_kv = create_inline_placeholders(excel_file)
    kv.update(inline_kv)

    wb = openpyxl.load_workbook(excel_file, data_only=True)
    list_placeholders = {}
    for sheet_name in wb.sheetnames:
        if sheet_name.startswith("By_"):
            key = sheet_name + "_List"
            items = build_list_from_sheet(excel_file, sheet_name)
            list_placeholders[key] = items

    toc_items = build_toc_from_sheet(excel_file, "Table_Contents")
    handle_toc_multi_slides(prs, toc_items)

    # Process all slides for replacements
    for slide in prs.slides:
        
        # *** ENHANCED: Use the new enhanced table processing function ***
        process_table_placeholders_with_expansion_enhanced(slide, list_placeholders, excel_file)
        
        # Regular bulleted lists (for text frames, not tables)
        for key, items in list_placeholders.items():
            placeholder = "{{" + key + "}}"
            if items:
                # Only process non-table placeholders here
                replace_list_placeholder_in_slide(slide, placeholder, items)

        # Text placeholders (includes inline keys)
        for key, val in kv.items():
            placeholder = "{{" + key + "}}"
            replace_text_placeholders_in_slide(slide, placeholder, val if val else "")

    # 🔽🔽🔽 NEW CODE BLOCK TO UPDATE CHARTS 🔽🔽🔽
    historical_years = list(range(2019, 2025))
    forecast_years = list(range(2025, 2034))
    for slide in prs.slides:
        update_charts_in_slide(slide, volumes, kv["Unit"], historical_years, forecast_years)
    # 🔼🔼🔼 END OF NEW BLOCK 🔼🔼🔼

    # Company table placeholders
    company_items = build_list_from_sheet(excel_file, "Company_Name")
    distribute_company_names_across_template_slides(prs, "{{Company_Name_List}}", company_items, duplicate_if_needed=True)

    prs.save(output_ppt)
    print("Saved:", output_ppt)

if __name__ == "__main__":
    import sys
    if len(sys.argv) == 4:
        excel_file, ppt_template, output_ppt = sys.argv[1], sys.argv[2], sys.argv[3]
    else:
        # fallback for local testing
        excel_file = "Datasheet-HS.xlsx"
        ppt_template = "default_template.pptx"
        output_ppt = "updated_presentation.pptx"

    main(excel_file, ppt_template, output_ppt)